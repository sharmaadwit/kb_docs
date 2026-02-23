import sys
from pptx import Presentation

def extract_text(pptx_path, md_path):
    print(f"Reading {pptx_path}")
    prs = Presentation(pptx_path)
    
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(f"# Presentation: {pptx_path.split('/')[-1]}\n\n")
        
        for i, slide in enumerate(prs.slides):
            f.write(f"## Slide {i+1}\n\n")
            
            # Check for title
            if slide.shapes.title and slide.shapes.title.text.strip():
                f.write(f"### {slide.shapes.title.text.strip()}\n\n")
            
            for shape in slide.shapes:
                # Skip title as we already added it
                if shape == slide.shapes.title:
                    continue
                    
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.replace('\v', '\n').replace('\r', '\n')
                    f.write(text + "\n\n")
                    
                # Optionally handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            text = cell.text.replace('\n', ' ').replace('\v', ' ').strip()
                            row_data.append(text)
                        f.write("| " + " | ".join(row_data) + " |\n")
                    f.write("\n")
                    
            f.write("---\n\n")
    print(f"Saved to {md_path}")

if __name__ == "__main__":
    extract_text(sys.argv[1], sys.argv[2])
